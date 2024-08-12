import os
import argparse
import time
import requests
import re
from pptx import Presentation
from pptx.util import Pt
from tqdm import tqdm
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Retrieve the OpenAI API key from the environment variable
API_KEY = os.getenv('OPENAI_API_KEY')

# Check for API key
if not API_KEY:
    raise ValueError("No API key found. Please set the 'OPENAI_API_KEY' environment variable.")

# Rate limiting configuration
'''
Here you need to find your rate limit and put it in the 2 variables below. Depending on if it is used anywhere else,
a safe starting point would be 50% of the max allowed per minute.

As I have a Tier 1 account, and I am not using my API anywhere else, I have just put the max for both according to api documentation found below:
https://platform.openai.com/docs/guides/rate-limits/usage-tiers?context=tier-one
'''
REQUEST_LIMIT = 3500  # Max requests per minute allowed
TOKEN_LIMIT = 200000  # Max tokens per minute allowed
REQUEST_COUNT = 0
TOKEN_COUNT = 0
START_TIME = time.time()

# Function to make sure we do not hit our rate limit 
def check_rate_limit(tokens):
    global REQUEST_COUNT, TOKEN_COUNT, START_TIME
    REQUEST_COUNT += 1
    TOKEN_COUNT += tokens

    elapsed_time = time.time() - START_TIME

    # If a minute has passed, reset the counters
    if elapsed_time >= 60:
        REQUEST_COUNT = 0
        TOKEN_COUNT = 0
        START_TIME = time.time()
        return

    # Check if we are exceeding the request or token limit
    if REQUEST_COUNT >= REQUEST_LIMIT or TOKEN_COUNT >= TOKEN_LIMIT:
        wait_time = 60 - elapsed_time
        print(f"Rate limit reached. Waiting for {wait_time:.2f} seconds...")
        time.sleep(wait_time)
        # Reset after waiting
        REQUEST_COUNT = 0
        TOKEN_COUNT = 0
        START_TIME = time.time()

# Function to estimate token count
def num_tokens(text):
    return len(text.split())

def contains_meaningful_content(text):
    """Check if the text contains any letters or numbers."""
    return bool(re.search(r'[a-zA-Z0-9]', text))

# POST translate text using OpenAI API
def translate_text(text, target_language):
    if not text.strip() or not contains_meaningful_content(text):
        return text
    
    tokens_needed = num_tokens(text)
    check_rate_limit(tokens_needed)  # Enforce rate limit before making the request
    
    url = "https://api.openai.com/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }
    
    system_message = f"""You are a translator specializing in PowerPoint presentations. 
    Your task is to translate text to {target_language}. 
    For image attributions or license information, keep proper nouns, abbreviations, and license codes unchanged. 
    Translate only the surrounding text.
    If the text looks like it should not be translated, then leave it as is (such as dates, math, equations, etc.).
    IMPORTANT: Your response must be in the following format:
    [START_TRANSLATION]
    Your translated text here
    [END_TRANSLATION]
    Any explanations or notes should be outside these tags."""

    user_message = f"Translate the following text to {target_language}:\n\n{text}"
    
    body = {
        "model": "gpt-3.5-turbo",
        "messages": [
            {"role": "system", "content": system_message},
            {"role": "user", "content": user_message}
        ],
        "max_tokens": 1000,
        "n": 1,
        "temperature": 0.3
    }
    
    response = requests.post(url, headers=headers, json=body)
    if response.status_code == 200:
        content = response.json()['choices'][0]['message']['content'].strip()
        start_tag = "[START_TRANSLATION]"
        end_tag = "[END_TRANSLATION]"
        start_index = content.find(start_tag) + len(start_tag)
        end_index = content.find(end_tag)
        if start_index != -1 and end_index != -1:
            return content[start_index:end_index].strip()
    return text

# Function to adjust font size to prevent overflow
def adjust_font_size(run, original_text, translated_text):
    original_length = len(original_text)
    translated_length = len(translated_text)
    
    if run.font.size is not None:
        current_font_size = run.font.size.pt
        if translated_length > original_length:
            scale_factor = original_length / translated_length
            new_font_size = current_font_size * scale_factor
        else:
            scale_factor = translated_length / original_length
            new_font_size = current_font_size * scale_factor
        
        # Ensure the new font size is within the valid range for PowerPoint
        new_font_size = max(10, min(new_font_size, 400))
        run.font.size = Pt(new_font_size)

# Translate the text within a shape
def translate_shape_text(shape, target_language):
    if hasattr(shape, "text_frame") and shape.text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip() and contains_meaningful_content(run.text):
                    original_text = run.text
                    translated_text = translate_text(run.text, target_language)
                    if translated_text != original_text:
                        adjust_font_size(run, original_text, translated_text)
                        run.text = translated_text

# Translate the text within a table
def translate_table(table, target_language):
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip() and contains_meaningful_content(run.text):
                        original_text = run.text
                        translated_text = translate_text(run.text, target_language)
                        if translated_text != original_text:
                            adjust_font_size(run, original_text, translated_text)
                            run.text = translated_text

# Flattened recursive function to process all shapes
def process_shapes_recursive(shapes, target_language):
    for shape in shapes:
        if shape.has_text_frame:
            translate_shape_text(shape, target_language)
        elif shape.has_table:
            translate_table(shape.table, target_language)
        elif hasattr(shape, 'shapes'):  # Check if it's a group shape
            process_shapes_recursive(shape.shapes, target_language)

# Process the entire presentation
def process_presentation(input_file, target_language):
    print(f"Opening {input_file}")
    try:
        input_ppt = Presentation(input_file)
    except Exception as e:
        print(f"Error opening file {input_file}: {e}")
        return

    slide_count = len(input_ppt.slides)
    
    with tqdm(total=slide_count, desc="Translating", unit="slide") as pbar:
        for slide in input_ppt.slides:
            process_shapes_recursive(slide.shapes, target_language)
            pbar.update(1)

    output_file = f"{target_language}_{os.path.basename(input_file)}"
    try:
        input_ppt.save(output_file)
        print(f"Saved as {output_file}")
    except Exception as e:
        print(f"Error saving file {output_file}: {e}")

# Process all presentations in a folder
def process_folder(folder_path, target_language):
    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx"):
            file_path = os.path.join(folder_path, filename)
            process_presentation(file_path, target_language)

# Main function to parse arguments and initiate processing
def main():
    parser = argparse.ArgumentParser(description="Translate PowerPoint presentations. Usage: python3 translatePPTX.py <input_path> <target_language>")
    parser.add_argument("input_path", nargs='?', help="Path to the input PowerPoint file or folder")
    parser.add_argument("target_language", nargs='?', help="Target language for translation (ie. Spanish, German, etc.)")
    args = parser.parse_args()

    if not args.input_path or not args.target_language:
        parser.print_help()
        return
    
    # handle individual vs multiple file handling
    if os.path.isdir(args.input_path):
        process_folder(args.input_path, args.target_language)
    else:
        process_presentation(args.input_path, args.target_language)

if __name__ == "__main__":
    main()
